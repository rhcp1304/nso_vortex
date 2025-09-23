import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.shapes.autoshape import Shape
from pptx.parts.image import ImagePart


def extract_data_from_ppt(filepath):
    """
    Extracts text, tables, and images from a PowerPoint presentation (.pptx file).

    Args:
        filepath (str): The path to the PowerPoint file.

    Returns:
        dict: A dictionary containing all extracted data, structured by slide.
              The keys are 'slide_data', and the value is a list of dictionaries,
              where each dictionary represents a slide's content.
    """
    if not os.path.exists(filepath):
        print(f"Error: The file '{filepath}' was not found.")
        return None

    try:
        prs = Presentation(filepath)
    except Exception as e:
        print(f"Error: Could not open the file. Is it a valid .pptx file? Error: {e}")
        return None

    presentation_data = {"slides": []}
    image_dir = "extracted_images"
    os.makedirs(image_dir, exist_ok=True)
    print(f"Images will be saved in the '{image_dir}' directory.")

    for slide_number, slide in enumerate(prs.slides):
        slide_content = {
            "slide_number": slide_number + 1,
            "title": "",
            "text": [],
            "tables": [],
            "images": []
        }

        # Iterate through all shapes on the slide
        for shape in slide.shapes:
            # Check for a title placeholder on the slide
            if shape.is_placeholder and shape.has_text_frame and shape.placeholder_format.type == 1:
                slide_content["title"] = shape.text_frame.text.strip()

            # --- Extract Text ---
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text.strip():
                        slide_content["text"].append(paragraph.text.strip())

            # --- Extract Tables ---
            if shape.has_table:
                table_data = []
                table = shape.table
                for row_idx, row in enumerate(table.rows):
                    row_data = []
                    for cell_idx, cell in enumerate(row.cells):
                        cell_text = cell.text.strip()
                        row_data.append(cell_text)
                    table_data.append(row_data)
                slide_content["tables"].append(table_data)

            # --- Extract Images ---
            if hasattr(shape, 'image'):
                image_bytes = shape.image.blob
                image_ext = shape.image.ext
                image_filename = f"slide_{slide_number + 1}_{shape.name}.{image_ext}"
                image_path = os.path.join(image_dir, image_filename)

                with open(image_path, "wb") as f:
                    f.write(image_bytes)
                slide_content["images"].append(image_path)

        presentation_data["slides"].append(slide_content)

    return presentation_data


def print_extracted_data(data):
    """
    Prints the extracted data from the PowerPoint file in a readable format.
    """
    if not data:
        return

    for slide in data["slides"]:
        print("\n" + "=" * 50)
        print(f"Slide {slide['slide_number']}:")
        if slide['title']:
            print(f"Title: {slide['title']}")

        if slide['text']:
            print("\nText:")
            for text_content in slide['text']:
                print(f"- {text_content}")

        if slide['tables']:
            print("\nTables:")
            for table in slide['tables']:
                for row in table:
                    print(f"  {row}")

        if slide['images']:
            print("\nImages:")
            for image_path in slide['images']:
                print(f"- Saved image: {image_path}")
